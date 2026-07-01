from pptx import Presentation

from app.ingestion.models.document import (
    Document
)

from app.ingestion.parsers.base_parser import (
    BaseParser
)


class PPTXParser(BaseParser):

    def parse(
        self,
        file_path: str
    ) -> Document:

        prs = Presentation(
            file_path
        )

        text = ""

        for slide in prs.slides:

            for shape in slide.shapes:

                if hasattr(
                    shape,
                    "text"
                ):
                    text += shape.text + "\n"

        return Document(
            content=text,
            metadata={
                "slides":
                len(prs.slides)
            },
            source=file_path,
            document_type="pptx"
        )